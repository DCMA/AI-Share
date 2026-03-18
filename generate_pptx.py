#!/usr/bin/env python3
"""
Generate a PowerPoint file from slides.json
Usage: python generate_pptx.py [slides.json] [output.pptx]
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Theme Colors ───
THEMES = {
    "dark": {
        "bg": RGBColor(0x0A, 0x0A, 0x0A),
        "surface": RGBColor(0x14, 0x14, 0x14),
        "border": RGBColor(0x2A, 0x2A, 0x2A),
        "text": RGBColor(0xFA, 0xFA, 0xFA),
        "text_secondary": RGBColor(0x88, 0x88, 0x88),
        "text_dim": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0xE0, 0xE0, 0xE0),
        "gradient_1": RGBColor(0xE0, 0xE0, 0xE0),
        "gradient_2": RGBColor(0xA0, 0xA0, 0xA0),
    },
    "midnight": {
        "bg": RGBColor(0x0B, 0x0D, 0x17),
        "surface": RGBColor(0x12, 0x15, 0x2A),
        "border": RGBColor(0x25, 0x2A, 0x4A),
        "text": RGBColor(0xE8, 0xEA, 0xF6),
        "text_secondary": RGBColor(0x79, 0x86, 0xCB),
        "text_dim": RGBColor(0x3F, 0x45, 0x80),
        "accent": RGBColor(0xB3, 0x88, 0xFF),
        "gradient_1": RGBColor(0xB3, 0x88, 0xFF),
        "gradient_2": RGBColor(0x53, 0x6D, 0xFE),
    },
    "warm": {
        "bg": RGBColor(0x1A, 0x14, 0x10),
        "surface": RGBColor(0x23, 0x1C, 0x16),
        "border": RGBColor(0x3D, 0x32, 0x28),
        "text": RGBColor(0xF5, 0xEB, 0xE0),
        "text_secondary": RGBColor(0xA0, 0x8A, 0x74),
        "text_dim": RGBColor(0x5A, 0x4A, 0x3A),
        "accent": RGBColor(0xF5, 0xC8, 0x7A),
        "gradient_1": RGBColor(0xF5, 0xC8, 0x7A),
        "gradient_2": RGBColor(0xC4, 0x95, 0x6A),
    },
    "light": {
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "surface": RGBColor(0xFF, 0xFF, 0xFF),
        "border": RGBColor(0xE0, 0xE0, 0xE0),
        "text": RGBColor(0x11, 0x11, 0x11),
        "text_secondary": RGBColor(0x66, 0x66, 0x66),
        "text_dim": RGBColor(0x99, 0x99, 0x99),
        "accent": RGBColor(0x11, 0x11, 0x11),
        "gradient_1": RGBColor(0x33, 0x33, 0x33),
        "gradient_2": RGBColor(0x88, 0x88, 0x88),
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_TITLE = "Inter"
FONT_MONO = "JetBrains Mono"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_TITLE, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_slide_number(slide, num, colors):
    add_textbox(
        slide,
        SLIDE_W - Inches(1.2), Inches(0.4), Inches(0.8), Inches(0.4),
        f"{num:02d}",
        font_size=10, color=colors["text_secondary"],
        alignment=PP_ALIGN.RIGHT, font_name=FONT_MONO
    )


def add_divider(slide, top, colors, center_x=None):
    left = center_x - Inches(0.4) if center_x else (SLIDE_W - Inches(0.8)) / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.8), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["border"]
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_badge(slide, text, top, colors):
    w = Inches(max(2, len(text) * 0.12 + 0.6))
    left = (SLIDE_W - w) / 2
    shape = add_rounded_rect(slide, left, top, w, Inches(0.35), colors["bg"], colors["border"])
    shape.adjustments[0] = 0.5  # pill shape
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = FONT_MONO
    p.font.color.rgb = colors["text_secondary"]
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_label(slide, text, left, top, colors):
    add_textbox(
        slide, left, top, Inches(4), Inches(0.3),
        text.upper(), font_size=9, color=colors["text_secondary"],
        font_name=FONT_MONO, bold=False
    )


# ═══════════════════════════════════════════════════════
# SLIDE TYPE RENDERERS
# ═══════════════════════════════════════════════════════

def render_cover(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, colors["bg"])

    y = Inches(1.8)
    if s.get("badge"):
        add_badge(slide, s["badge"], y, colors)
        y += Inches(0.7)

    title = s.get("title", "").replace("\\n", "\n")
    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(2),
        title, font_size=48, color=colors["accent"], bold=True,
        alignment=PP_ALIGN.CENTER, line_spacing=1.05
    )
    y += Inches(2.2)

    if s.get("subtitle"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.6),
            s["subtitle"], font_size=20, color=colors["text_secondary"],
            alignment=PP_ALIGN.CENTER
        )
        y += Inches(0.8)

    add_divider(slide, y, colors)
    y += Inches(0.5)

    if s.get("author"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.4),
            s["author"], font_size=11, color=colors["text_dim"],
            alignment=PP_ALIGN.CENTER
        )

    add_slide_number(slide, idx + 1, colors)


def render_end(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(2.2)
    if s.get("badge"):
        add_badge(slide, s["badge"], y, colors)
        y += Inches(0.7)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(1.5),
        s.get("title", "Thank You"), font_size=48, color=colors["accent"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    y += Inches(1.5)

    if s.get("subtitle"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.5),
            s["subtitle"], font_size=20, color=colors["text_secondary"],
            alignment=PP_ALIGN.CENTER
        )
        y += Inches(0.7)

    add_divider(slide, y, colors)
    y += Inches(0.5)

    author_text = s.get("author", "")
    if s.get("handle"):
        author_text += f" · {s['handle']}"
    if author_text:
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.4),
            author_text, font_size=11, color=colors["text_dim"],
            alignment=PP_ALIGN.CENTER
        )

    add_slide_number(slide, idx + 1, colors)


def render_section(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(2.5)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(2),
        s.get("title", "").replace("\\n", "\n"), font_size=44, color=colors["accent"],
        bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.1
    )
    y += Inches(2)

    if s.get("subtitle"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.6),
            s["subtitle"], font_size=18, color=colors["text_secondary"],
            alignment=PP_ALIGN.CENTER
        )

    add_slide_number(slide, idx + 1, colors)


def render_list(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    is_left = s.get("align") == "left"
    x = Inches(1.5) if is_left else (SLIDE_W - Inches(8)) / 2
    y = Inches(1.2)
    align = PP_ALIGN.LEFT if is_left else PP_ALIGN.CENTER

    if s.get("label"):
        add_label(slide, s["label"], x, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, x, y, Inches(8), Inches(1),
        s.get("title", "").replace("\\n", "\n"), font_size=36,
        color=colors["text"], bold=True, alignment=align
    )
    y += Inches(1.2)

    for item in s.get("items", []):
        # bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.18), Pt(6), Pt(6)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors["text_secondary"]
        dot.line.fill.background()
        dot.shadow.inherit = False

        add_textbox(
            slide, x + Inches(0.5), y, Inches(7), Inches(0.5),
            item, font_size=15, color=colors["text_secondary"]
        )

        # separator line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y + Inches(0.55), Inches(8), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["border"]
        line.line.fill.background()
        line.shadow.inherit = False

        y += Inches(0.65)

    add_slide_number(slide, idx + 1, colors)


def render_cards(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(0.8)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(1),
        s.get("title", "").replace("\\n", "\n"), font_size=34,
        color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.1
    )
    y += Inches(1.3)

    cards = s.get("cards", [])
    cols = min(s.get("columns", 3), len(cards))
    card_w = Inches(3.2)
    gap = Inches(0.3)
    total_w = card_w * cols + gap * (cols - 1)
    start_x = (SLIDE_W - total_w) / 2

    for i, card in enumerate(cards):
        cx = start_x + (card_w + gap) * (i % cols)
        cy = y + (Inches(2.8)) * (i // cols)

        rect = add_rounded_rect(slide, cx, cy, card_w, Inches(2.5), colors["surface"], colors["border"])

        icon_text = card.get("icon", "")
        add_textbox(slide, cx + Inches(0.3), cy + Inches(0.3), Inches(1), Inches(0.4),
                    icon_text, font_size=16, color=colors["text_dim"], bold=True)

        add_textbox(slide, cx + Inches(0.3), cy + Inches(0.8), card_w - Inches(0.6), Inches(0.4),
                    card["title"], font_size=14, color=colors["text"], bold=True)

        add_textbox(slide, cx + Inches(0.3), cy + Inches(1.3), card_w - Inches(0.6), Inches(1),
                    card["desc"], font_size=12, color=colors["text_secondary"], line_spacing=1.5)

    add_slide_number(slide, idx + 1, colors)


def render_stats(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(1.5)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(1),
        s.get("title", "").replace("\\n", "\n"), font_size=34,
        color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER
    )
    y += Inches(1.5)

    stats = s.get("stats", [])
    n = len(stats)
    stat_w = Inches(3)
    total_w = stat_w * n
    start_x = (SLIDE_W - total_w) / 2

    for i, st in enumerate(stats):
        sx = start_x + stat_w * i
        add_textbox(
            slide, sx, y, stat_w, Inches(1),
            st["number"], font_size=48, color=colors["accent"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, sx, y + Inches(1), stat_w, Inches(0.4),
            st["desc"], font_size=13, color=colors["text_secondary"],
            alignment=PP_ALIGN.CENTER
        )

    add_slide_number(slide, idx + 1, colors)


def render_flow(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(1.2)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(1),
        s.get("title", "").replace("\\n", "\n"), font_size=34,
        color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER
    )
    y += Inches(1.0)

    if s.get("subtitle"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.6),
            s["subtitle"], font_size=16, color=colors["text_secondary"],
            alignment=PP_ALIGN.CENTER
        )
        y += Inches(0.9)

    steps = s.get("steps", [])
    n = len(steps)
    step_w = Inches(1.6)
    arrow_w = Inches(0.5)
    total_w = step_w * n + arrow_w * (n - 1)
    start_x = (SLIDE_W - total_w) / 2

    for i, st in enumerate(steps):
        sx = start_x + (step_w + arrow_w) * i
        rect = add_rounded_rect(slide, sx, y, step_w, Inches(1.2), colors["surface"], colors["border"])
        rect.adjustments[0] = 0.08

        add_textbox(slide, sx + Inches(0.1), y + Inches(0.15), step_w - Inches(0.2), Inches(0.4),
                    st["name"], font_size=12, color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx + Inches(0.1), y + Inches(0.55), step_w - Inches(0.2), Inches(0.4),
                    st["desc"], font_size=9, color=colors["text_dim"], alignment=PP_ALIGN.CENTER)

        if i < n - 1:
            ax = sx + step_w + Inches(0.05)
            add_textbox(slide, ax, y + Inches(0.3), arrow_w, Inches(0.4),
                        "→", font_size=16, color=colors["text_dim"], alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, idx + 1, colors)


def render_steps(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    is_left = s.get("align") == "left"
    x = Inches(1.5) if is_left else Inches(2)
    y = Inches(0.8)
    align = PP_ALIGN.LEFT if is_left else PP_ALIGN.CENTER

    if s.get("label"):
        add_label(slide, s["label"], x, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, x, y, Inches(8), Inches(0.8),
        s.get("title", "").replace("\\n", "\n"), font_size=32,
        color=colors["text"], bold=True, alignment=align
    )
    y += Inches(1.0)

    card_w = Inches(9)
    card_h = Inches(1.1)

    for i, step in enumerate(s.get("steps", [])):
        rect = add_rounded_rect(slide, x, y, card_w, card_h, colors["surface"], colors["border"])

        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.6),
                    f"{i+1:02d}", font_size=22, color=colors["text_dim"], bold=True)

        add_textbox(slide, x + Inches(1.2), y + Inches(0.15), Inches(7), Inches(0.4),
                    step["title"], font_size=14, color=colors["text"], bold=True)

        add_textbox(slide, x + Inches(1.2), y + Inches(0.55), Inches(7), Inches(0.4),
                    step["desc"], font_size=11, color=colors["text_secondary"])

        y += card_h + Inches(0.15)

    add_slide_number(slide, idx + 1, colors)


def render_code(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(1.5)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    add_textbox(
        slide, Inches(1), y, SLIDE_W - Inches(2), Inches(1),
        s.get("title", "").replace("\\n", "\n"), font_size=34,
        color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER
    )
    y += Inches(1.2)

    # Code block background
    code_w = Inches(8)
    code_h = Inches(3.5)
    code_x = (SLIDE_W - code_w) / 2
    rect = add_rounded_rect(slide, code_x, y, code_w, code_h, colors["surface"], colors["border"])
    rect.adjustments[0] = 0.03

    # Build code text
    code_lines = []
    for c in s.get("code", []):
        if c["type"] == "break":
            code_lines.append("")
        else:
            code_lines.append(c.get("text", ""))
    code_text = "\n".join(code_lines)

    add_textbox(
        slide, code_x + Inches(0.4), y + Inches(0.3), code_w - Inches(0.8), code_h - Inches(0.6),
        code_text, font_size=12, color=colors["text_secondary"],
        font_name=FONT_MONO, line_spacing=1.6
    )

    add_slide_number(slide, idx + 1, colors)


def render_quote(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(2.0)
    if s.get("badge"):
        add_badge(slide, s["badge"], y, colors)
        y += Inches(0.6)

    add_divider(slide, y, colors)
    y += Inches(0.6)

    quote_text = '"' + s.get("text", "").replace("\\n", "\n") + '"'
    add_textbox(
        slide, Inches(2), y, SLIDE_W - Inches(4), Inches(2),
        quote_text, font_size=22, color=colors["text_secondary"],
        alignment=PP_ALIGN.CENTER, line_spacing=1.5
    )
    y += Inches(2.2)

    if s.get("author"):
        add_textbox(
            slide, Inches(2), y, SLIDE_W - Inches(4), Inches(0.4),
            f"— {s['author']}", font_size=11, color=colors["text_dim"],
            alignment=PP_ALIGN.CENTER
        )

    add_slide_number(slide, idx + 1, colors)


def render_image(prs, s, idx, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y = Inches(0.8)
    if s.get("label"):
        add_label(slide, s["label"], (SLIDE_W - Inches(4)) / 2, y, colors)
        y += Inches(0.5)

    if s.get("title"):
        add_textbox(
            slide, Inches(1), y, SLIDE_W - Inches(2), Inches(0.8),
            s["title"].replace("\\n", "\n"), font_size=30,
            color=colors["text"], bold=True, alignment=PP_ALIGN.CENTER
        )
        y += Inches(1.0)

    if s.get("src"):
        try:
            img_w = Inches(8)
            slide.shapes.add_picture(s["src"], (SLIDE_W - img_w) / 2, y, img_w)
        except Exception:
            add_textbox(
                slide, Inches(2), y, SLIDE_W - Inches(4), Inches(1),
                f"[Image: {s['src']}]", font_size=14, color=colors["text_dim"],
                alignment=PP_ALIGN.CENTER
            )

    add_slide_number(slide, idx + 1, colors)


# ═══════════════════════════════════════════════════════

RENDERERS = {
    "cover": render_cover,
    "end": render_end,
    "section": render_section,
    "list": render_list,
    "cards": render_cards,
    "stats": render_stats,
    "flow": render_flow,
    "steps": render_steps,
    "code": render_code,
    "quote": render_quote,
    "image": render_image,
}


def main():
    json_path = sys.argv[1] if len(sys.argv) > 1 else "slides.json"
    output_path = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    theme_name = data.get("theme", {}).get("preset", "dark")
    colors = THEMES.get(theme_name, THEMES["dark"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, s in enumerate(data.get("slides", [])):
        renderer = RENDERERS.get(s.get("type"), render_section)
        renderer(prs, s, i, colors)

    prs.save(output_path)
    print(f"✓ Generated: {output_path}  ({len(data['slides'])} slides, theme: {theme_name})")


if __name__ == "__main__":
    main()
